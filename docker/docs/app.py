"""aichat-docs: Pure-text document ingestor. No ML — outputs normalized Markdown.

Supported formats: .pdf, .docx, .xlsx, .pptx, .html, .md, .txt

Endpoints:
  POST /ingest      — convert document (base64) to Markdown
  POST /ingest/url  — download URL and ingest
  POST /tables      — extract tables from document
  GET  /formats     — list supported extensions
  GET  /health
"""
from __future__ import annotations

import asyncio
import base64
import io
import logging
import os
import re
from pathlib import Path

import httpx
from fastapi import FastAPI, HTTPException, Request
from fastapi.responses import JSONResponse

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
log = logging.getLogger("aichat-docs")

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

DB_API   = os.environ.get("DATABASE_URL", "http://aichat-database:8091")
_SERVICE = "aichat-docs"

SUPPORTED = ["pdf", "docx", "xlsx", "xls", "pptx", "html", "htm", "md", "txt"]

app = FastAPI()


# ---------------------------------------------------------------------------
# Error reporting
# ---------------------------------------------------------------------------

async def _report_error(message: str, detail: str | None = None) -> None:
    try:
        async with httpx.AsyncClient(timeout=5) as client:
            await client.post(
                f"{DB_API}/errors/log",
                json={"service": _SERVICE, "level": "ERROR",
                      "message": message, "detail": detail},
            )
    except Exception:
        pass


@app.exception_handler(Exception)
async def _global_exc(request: Request, exc: Exception) -> JSONResponse:
    msg = str(exc)
    log.error("Unhandled [%s %s]: %s", request.method, request.url.path, msg, exc_info=True)
    asyncio.create_task(_report_error(msg, f"{request.method} {request.url.path}"))
    return JSONResponse(status_code=500, content={"error": msg})


# ---------------------------------------------------------------------------
# Extraction helpers
# ---------------------------------------------------------------------------

def _extract_pdf(data: bytes) -> tuple[str, list[dict]]:
    """Extract text and tables from PDF using pdfminer."""
    from pdfminer.high_level import extract_text
    from pdfminer.layout import LAParams
    text = extract_text(io.BytesIO(data), laparams=LAParams())
    # pdfminer doesn't have a great table extractor — return text only
    pages = text.split("\x0c")
    md_parts = []
    for i, page in enumerate(pages):
        stripped = page.strip()
        if stripped:
            md_parts.append(f"## Page {i + 1}\n\n{stripped}")
    return "\n\n".join(md_parts), []


def _extract_docx(data: bytes) -> tuple[str, list[dict]]:
    """Extract text and tables from .docx."""
    import docx
    doc  = docx.Document(io.BytesIO(data))
    parts = []
    tables = []

    # Paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style:
            parts.append(f"# {text}")
        elif "Heading 2" in style:
            parts.append(f"## {text}")
        elif "Heading 3" in style:
            parts.append(f"### {text}")
        else:
            parts.append(text)

    # Tables
    for i, table in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        headers = rows[0]
        body    = rows[1:]
        tbl = {"title": f"Table {i + 1}", "headers": headers, "rows": body}
        tables.append(tbl)
        # Render as markdown table
        sep  = " | ".join("---" for _ in headers)
        head = " | ".join(headers)
        body_lines = [" | ".join(r) for r in body]
        parts.append(f"\n{head}\n{sep}\n" + "\n".join(body_lines))

    return "\n\n".join(parts), tables


def _extract_xlsx(data: bytes) -> tuple[str, list[dict]]:
    """Extract sheets as Markdown tables from .xlsx."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    tables = []
    for sheet in wb.worksheets:
        rows = [[str(cell.value or "") for cell in row] for row in sheet.iter_rows()]
        if not rows:
            continue
        parts.append(f"## Sheet: {sheet.title}")
        headers = rows[0]
        body    = rows[1:]
        tbl = {"title": sheet.title, "headers": headers, "rows": body}
        tables.append(tbl)
        sep  = " | ".join("---" for _ in headers)
        head = " | ".join(headers)
        body_lines = [" | ".join(r) for r in body[:200]]  # cap at 200 rows in output
        parts.append(f"{head}\n{sep}\n" + "\n".join(body_lines))
    return "\n\n".join(parts), tables


def _extract_pptx(data: bytes) -> tuple[str, list[dict]]:
    """Extract slide text from .pptx."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = [f"## Slide {i + 1}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        parts.append("\n\n".join(slide_parts))
    return "\n\n".join(parts), []


def _extract_html(data: bytes) -> tuple[str, list[dict]]:
    """Extract text and tables from HTML."""
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(data, "html.parser")
    # Remove scripts/styles
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()

    tables = []
    for i, tbl in enumerate(soup.find_all("table")):
        rows = []
        for tr in tbl.find_all("tr"):
            cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
            if cells:
                rows.append(cells)
        if rows:
            tables.append({"title": f"Table {i+1}", "headers": rows[0], "rows": rows[1:]})
        tbl.replace_with(f"\n[Table {i+1}]\n")

    text = soup.get_text(separator="\n", strip=True)
    # Collapse blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text, tables


def _extract_text(data: bytes) -> tuple[str, list[dict]]:
    return data.decode("utf-8", errors="replace"), []


def _ingest(data: bytes, filename: str) -> dict:
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext == "pdf":
        md, tables = _extract_pdf(data)
    elif ext == "docx":
        md, tables = _extract_docx(data)
    elif ext in ("xlsx", "xls"):
        md, tables = _extract_xlsx(data)
    elif ext == "pptx":
        md, tables = _extract_pptx(data)
    elif ext in ("html", "htm"):
        md, tables = _extract_html(data)
    elif ext in ("md", "txt", ""):
        md, tables = _extract_text(data)
    else:
        raise HTTPException(status_code=422,
                            detail=f"Unsupported format '.{ext}'. Use: {SUPPORTED}")

    # Derive title from first heading or first line
    title = ""
    for line in md.splitlines():
        line = line.strip().lstrip("#").strip()
        if line:
            title = line[:120]
            break

    words = md.split()
    return {
        "markdown":    md,
        "title":       title,
        "word_count":  len(words),
        "tables_found": len(tables),
        "tables":       tables,
        "format":      ext or "txt",
    }


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@app.get("/health")
def health() -> dict:
    return {"status": "ok", "formats": SUPPORTED}


@app.get("/formats")
def formats() -> dict:
    return {"formats": SUPPORTED}


@app.post("/ingest")
def ingest_b64(payload: dict) -> dict:
    b64      = str(payload.get("b64", "")).strip()
    filename = str(payload.get("filename", "document.txt")).strip() or "document.txt"
    if not b64:
        raise HTTPException(status_code=422, detail="'b64' is required")
    try:
        data = base64.b64decode(b64)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Invalid base64: {exc}") from exc
    result = _ingest(data, filename)
    log.info("ingest %s → %d words, %d tables", filename, result["word_count"], result["tables_found"])
    return result


@app.post("/ingest/url")
async def ingest_url(payload: dict) -> dict:
    url = str(payload.get("url", "")).strip()
    if not url:
        raise HTTPException(status_code=422, detail="'url' is required")
    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=True, verify=False) as client:
            resp = await client.get(url)
            resp.raise_for_status()
            data = resp.content
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Failed to fetch URL: {exc}") from exc

    # Detect filename from URL or Content-Type
    filename = url.split("?")[0].rstrip("/").split("/")[-1] or "document"
    content_type = resp.headers.get("content-type", "")
    ext_candidate = Path(filename).suffix.lower().lstrip(".")
    if "." not in filename or ext_candidate not in SUPPORTED:
        if "pdf" in content_type:
            filename = "document.pdf"
        elif "html" in content_type:
            filename = "document.html"
        else:
            filename = "document.txt"

    result = _ingest(data, filename)
    result["url"] = url
    log.info("ingest_url %s → %d words", url, result["word_count"])
    return result


@app.post("/tables")
def extract_tables_only(payload: dict) -> dict:
    b64      = str(payload.get("b64", "")).strip()
    filename = str(payload.get("filename", "document.txt")).strip() or "document.txt"
    if not b64:
        raise HTTPException(status_code=422, detail="'b64' is required")
    try:
        data = base64.b64decode(b64)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Invalid base64: {exc}") from exc
    result = _ingest(data, filename)
    return {"tables": result["tables"], "count": len(result["tables"])}
